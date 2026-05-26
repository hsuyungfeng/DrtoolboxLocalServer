import os
import glob
import shutil
from config.settings import SPECIAL_DATA_DIR, GENERAL_DATA_DIR, DATA_DIR
import logging

logger = logging.getLogger(__name__)

# Setup archive directory
ARCHIVE_DIR = os.path.join(DATA_DIR, "archive")
os.makedirs(ARCHIVE_DIR, exist_ok=True)

_whisper_model = None

def get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        from faster_whisper import WhisperModel
        import torch
        device = "cuda" if torch.cuda.is_available() else "cpu"
        compute_type = "float16" if device == "cuda" else "int8"
        logger.info(f"Loading Whisper model on {device} with {compute_type}...")
        _whisper_model = WhisperModel("small", device=device, compute_type=compute_type)
    return _whisper_model

def _do_pdf_ocr(pdf_path):
    from PIL import Image
    import pytesseract
    from pdf2image import convert_from_path
    try:
        # Lower DPI for memory safety, and limit pages processed to prevent freezing
        images = convert_from_path(pdf_path, dpi=150, last_page=20)
        text = ""
        for i, img in enumerate(images):
            # Preprocess image
            max_size = 3000
            if img.width > max_size or img.height > max_size:
                img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
            
            try:
                page_text = pytesseract.image_to_string(img.convert('L'), lang='chi_tra+eng')
            except Exception:
                page_text = pytesseract.image_to_string(img.convert('L'))
            text += f"--- Page {i+1} ---\n{page_text}\n"
        return text
    except Exception as e:
        logger.error(f"PDF OCR failed for {pdf_path}: {e}")
        return ""

def extract_text_from_file(filepath):
    from PIL import Image
    import pytesseract
    ext = filepath.lower().split('.')[-1]
    try:
        if ext in ['txt', 'md']:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == 'pdf':
            import PyPDF2
            text = ""
            try:
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages[:30]: # Limit reading to first 30 pages
                        extracted = page.extract_text()
                        if extracted:
                            text += extracted + "\n"
            except: pass
            
            if len(text.strip()) < 15:
                logger.info(f"PDF {filepath} 似乎是掃描檔或讀取失敗，啟動 OCR 備援機制...")
                text = _do_pdf_ocr(filepath)
            return text
        elif ext in ['jpg', 'jpeg', 'png']:
            try:
                with Image.open(filepath) as img:
                    # Resize if massive
                    if img.width > 4000 or img.height > 4000:
                        img.thumbnail((3000, 3000), Image.Resampling.LANCZOS)
                    # Grayscale for OCR
                    img = img.convert('L')
                    try:
                        return pytesseract.image_to_string(img, lang='chi_tra+eng')
                    except:
                        return pytesseract.image_to_string(img)
            except Exception as ocr_e:
                logger.error(f"Image OCR failed for {filepath}: {ocr_e}")
                return ""
        elif ext in ['mp4', 'mp3', 'm4a', 'wav', 'flv']:
            logger.info(f"Starting Whisper transcription for {filepath}")
            try:
                model = get_whisper_model()
                segments, info = model.transcribe(filepath, beam_size=5)
                segments_list = list(segments)
                if not segments_list:
                    return "--- 語音逐字稿 (未偵測到任何語音內容) ---"
                    
                text = f"--- 語音逐字稿 (語言: {info.language}) ---\n"
                for segment in segments_list:
                    text += f"[{segment.start:.2f}s -> {segment.end:.2f}s] {segment.text}\n"
                return text
            except Exception as whisper_e:
                logger.error(f"Whisper transcription failed for {filepath}: {whisper_e}")
                return f"--- 語音逐字稿失敗: {whisper_e} ---"

        # === Office File Parsing ===
        text = ""
        if ext == 'docx':
            import docx
            doc = docx.Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext in ['doc', 'ppt']:
            import subprocess, tempfile
            logger.info(f"Using LibreOffice to extract text from legacy {ext} file: {filepath}")
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    subprocess.run(['soffice', '--headless', '--convert-to', 'txt:Text', '--outdir', temp_dir, filepath], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                    base_name = os.path.splitext(os.path.basename(filepath))[0]
                    txt_file = os.path.join(temp_dir, f"{base_name}.txt")
                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                            text = f.read()
            except Exception as e:
                logger.error(f"Failed to convert legacy {ext} file {filepath}: {e}")
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        if ext in ['doc', 'docx', 'ppt', 'pptx'] and len(text.strip()) < 15:
            logger.info(f"{ext} 檔案 {filepath} 似乎是由純圖片組成，啟動 OCR 備援機制...")
            text = _do_pdf_ocr(filepath) # Re-use PDF OCR logic (works because we convert office to PDF)

        return text
    except Exception as e:
        logger.error(f"Error extracting text from {filepath}: {e}")
        return ""

def archive_file(filepath):
    """Moves a file to the archive folder instead of deleting it."""
    try:
        filename = os.path.basename(filepath)
        dest_path = os.path.join(ARCHIVE_DIR, filename)
        if os.path.exists(dest_path):
            import uuid
            dest_path = os.path.join(ARCHIVE_DIR, f"{uuid.uuid4().hex[:8]}_{filename}")
        shutil.move(filepath, dest_path)
        logger.info(f"[Safe Archive] Moved original to archive: {filename}")
        return True
    except Exception as e:
        logger.error(f"[Safe Archive] Failed to archive {filepath}: {e}")
        return False

def process_and_load_directory(directory, rag_engine_instance=None, is_special=True):
    documents = []
    if not os.path.exists(directory):
        return documents
        
    for filepath in glob.glob(os.path.join(directory, "**/*.txt"), recursive=True):
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
                if content.strip():
                    documents.append({"id": filepath, "content": content})
        except Exception as e:
            logger.error(f"Error reading {filepath}: {e}")
            
    if rag_engine_instance:
        import threading
        def background_ocr():
            supported_exts = ['.pdf', '.doc', '.docx', '.ppt', '.pptx', '.jpg', '.jpeg', '.png', '.mp4', '.mp3', '.m4a', '.wav', '.flv', '.md']
            for root, _, files in os.walk(directory):
                for file in files:
                    ext = os.path.splitext(file)[1].lower()
                    if ext in supported_exts:
                        filepath = os.path.join(root, file)
                        txt_path = filepath + ".txt"
                        alt_txt_path = os.path.splitext(filepath)[0] + ".txt"
                        if os.path.exists(txt_path) or os.path.exists(alt_txt_path):
                            archive_file(filepath)
                            continue
                        logger.info(f"[Background OCR] Processing: {file}")
                        extracted_text = extract_text_from_file(filepath)
                        if extracted_text and extracted_text.strip():
                            try:
                                with open(txt_path, 'w', encoding='utf-8') as f:
                                    f.write(extracted_text)
                                doc = {"id": txt_path, "content": extracted_text}
                                if is_special: rag_engine_instance.ingest_special_data([doc])
                                else: rag_engine_instance.ingest_general_data([doc])
                                archive_file(filepath)
                            except Exception as e:
                                logger.error(f"Failed to save extracted text for {filepath}: {e}")
        threading.Thread(target=background_ocr, daemon=True).start()
    return documents

def get_special_data(rag_engine_instance=None):
    logger.info(f"Loading special data from {SPECIAL_DATA_DIR}")
    return process_and_load_directory(SPECIAL_DATA_DIR, rag_engine_instance, is_special=True)

def get_general_data(rag_engine_instance=None):
    logger.info(f"Loading general data from {GENERAL_DATA_DIR}")
    return process_and_load_directory(GENERAL_DATA_DIR, rag_engine_instance, is_special=False)
